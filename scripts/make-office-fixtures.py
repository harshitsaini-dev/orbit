"""Regenerates the Office samples the reader tests run against.

Written by the libraries that write the real formats, not by hand: a fixture
built by the same person who wrote the reader only proves the two agree.

    pip install openpyxl python-docx python-pptx
    python scripts/make-office-fixtures.py
"""

import datetime
import os

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

OUT = os.path.join("apps", "web", "src", "lib", "__fixtures__")


def spreadsheet() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Units", "Revenue", "Signed"])
    ws.append(["North", 120, 4500.5, datetime.date(2026, 3, 14)])
    ws.append(["South", 80, 2200, datetime.date(2026, 4, 1)])
    # A gap: this row has A and D but nothing between, which is where a reader
    # that ignores cell references shifts every later column left.
    ws["A5"] = "East"
    ws["D5"] = datetime.date(2026, 5, 20)

    notes = wb.create_sheet("Notes")
    notes.append(["Quoted, with comma"])
    notes.append(['He said "hi"'])

    wb.save(os.path.join(OUT, "sample.xlsx"))


def document() -> None:
    doc = Document()
    doc.add_heading("Quarterly report", level=1)
    doc.add_paragraph("An ordinary paragraph with some text in it.")
    doc.add_heading("Findings", level=2)
    doc.add_paragraph("First finding", style="List Bullet")
    doc.add_paragraph("Second finding", style="List Bullet")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Units"
    table.cell(1, 1).text = "200"

    doc.add_paragraph("Closing paragraph after the table.")
    doc.save(os.path.join(OUT, "sample.docx"))


def presentation() -> None:
    prs = Presentation()

    first = prs.slides.add_slide(prs.slide_layouts[1])
    first.shapes.title.text = "Orbit"
    first.placeholders[1].text = "One workspace\nEvery cloud"
    first.notes_slide.notes_text_frame.text = "Remember to mention the proxy."

    second = prs.slides.add_slide(prs.slide_layouts[1])
    second.shapes.title.text = "Second slide"
    second.placeholders[1].text = "Just one point"

    prs.save(os.path.join(OUT, "sample.pptx"))


if __name__ == "__main__":
    os.makedirs(OUT, exist_ok=True)
    spreadsheet()
    document()
    presentation()
    for name in ("sample.xlsx", "sample.docx", "sample.pptx"):
        print(name, os.path.getsize(os.path.join(OUT, name)), "bytes")
